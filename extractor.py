from pathlib import Path


def extract_text(file_path: str) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    elif suffix == ".docx":
        return _extract_docx(path)
    elif suffix == ".pptx":
        return _extract_pptx(path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


def _extract_pdf(path):
    import pdfplumber
    text_parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
    if not text_parts:
        raise ValueError("No text could be extracted from this PDF. It may be a scanned image.")
    return "\n\n".join(text_parts)


def _extract_docx(path):
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract table content
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n".join(paragraphs)


def _extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    if not slides:
        raise ValueError("No text could be extracted from this PowerPoint.")
    return "\n\n".join(slides)
